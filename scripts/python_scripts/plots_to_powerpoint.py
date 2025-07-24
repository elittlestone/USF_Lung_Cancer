from pptx import Presentation
from pptx.util import Inches, Pt
from collections import defaultdict
import os
import re
import argparse 

parser = argparse.ArgumentParser(description = "Python script for outputting plots to Powerpoint")
parser.add_argument("--output_filepath", type = str, required = True, 
                    help = "Filepath for outputted powerpoint file")

parser.add_argument("--plots_location", type = str, required = True,
                    help = "Directory where plots are located")

args = parser.parse_args()

# Folder for plots to be stored in
plots_dir = args.plots_location

plot_files = [file for file in os.listdir(plots_dir) if not file.startswith("umap_celltypes")]
cell_type_umap_path = os.path.join(plots_dir, "umap_celltypes.png")

# Group plots by gene
gene_to_plots = defaultdict(dict)

for file in plot_files:
    match = re.match(r"(umap|violin|dotplot)_(.+?)_", file)
    if match:
        plot_type, gene = match.groups()
        gene_to_plots[gene][plot_type] = os.path.join(plots_dir, file)

# Define slide dimensions and create Presentation object 
pres = Presentation()
pres.slide_width = Inches(14)
pres.slide_height = Inches(16)
blank_slide_layout = pres.slide_layouts[6]
#image_size = Inches(10)

positions = [
        (Inches(0), Inches(0)),
        (Inches(9), Inches(0)),
        (Inches(0), Inches(9)),
        (Inches(9), Inches(9))
        ]

for gene, plots in gene_to_plots.items():
    slide = pres.slides.add_slide(blank_slide_layout)
    
    # Gather all plots to place in the slide 
    plot_images = []
    
    # For each slide/gene, include the celltypes umap 
    plot_images.append(cell_type_umap_path)

    # Add gene specific plots
    for plot_type in ["umap", "violin", "dotplot"]:
        if plot_type in plots:
            plot_images.append(plots[plot_type])

    # Put all 4 plots in 2x2 grid 
    for i, image_path in enumerate(plot_images[:]):
        left, top = positions[i]
        slide.shapes.add_picture(image_path, left, top)

pres.save(args.output_filepath)

